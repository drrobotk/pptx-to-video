 #!/usr/bin/env python

import os, argparse, re, fitz, shutil, comtypes.client
from progress.bar import Bar
from pptx import Presentation
from gtts import gTTS
from subprocess import DEVNULL, STDOUT, call

__author__ = ['Dr. Usman Kayani']

def pptx_video(
    pptx_file: str,
    output_file: str,
) -> None:
    """
    Convert a powerpoint presentation (pptx) to video (mp4) using 
    text-to-speech with the presenter notes or slide text.
    
    Parameters
    ----------
    pptx_file: str
        Input file path to pptx file to convert.
    output_file: str
        Output file path to mp4 file.
    
    Returns
    -------
    None
    """
    name = pptx_file[:-5]
    pdf_file = f'{name}.pdf'

    # Create PDF file for image processing, if it doesn't exist
    if not os.path.exists(pdf_file):
        _pptx_pdf(pptx_file)
        
    os.system('mkdir tmp')

    prs = Presentation(pptx_file)
    doc = fitz.open(pdf_file)
    N = len(prs.slides)

    # The number of pages of the PPTX file and PDF should be equal.
    assert doc.page_count == N

    print('Creating MP3s (TTS) and PNGs for each slide, then combining them into MP4...')
    bar = Bar('Processing slides', max=5*N, suffix='%(percent)d%%')
    for slide in prs.slides:
        i = prs.slides.index(slide)

        page = doc.load_page(i)
        speech = _speech_text(slide)
        bar.next()
        
        pix = page.get_pixmap()
        bar.next()

        # Get audio from text-to-speech on slide text.
        try:
            tts = gTTS(text=speech, lang='en')
        except:
            notes = ' '.join(re.split('\s+', notes, flags=re.UNICODE))
            tts = gTTS(text=speech, lang='en')
        bar.next()

        file_name = f'tmp\{name}_slide_{i+1}'
        image_path, audio_path, out_path_mp4 = _get_filepaths(file_name)
            
        pix.save(image_path)
        tts.save(audio_path)
        bar.next()
        ffmpeg_args = f'-y -loop 1 -i {image_path} -i {out_path_mp4} -c:v libx264 ' + \
        f'-tune stillimage -c:a aac -b:a 192k -pix_fmt yuv420p -shortest {output_file}'
        _execute_cmd(f'ffmpeg {ffmpeg_args}')
        bar.next()
        f = open('tmp\list.txt', 'w')
        f.write(f'file tmp/{name}_slide_{i+1}.mp4\n')
    f.close()

    print(f'\n Combining the MP4s for all slides into the single output video {output_file}...', end='')
    ffmpeg_args = f'-y -f concat -i tmp\list.txt -c copy {output_file}'
    _execute_cmd(f'ffmpeg {ffmpeg_args}')
    print('Done!')
    # shutil.rmtree('tmp')
    os.remove(pdf_file)
    # shutil.rmtree('tmp', ignore_errors=True)
    print('done!')
    bar.finish()

def _get_filepaths(
    file_name: str,
) -> str:
    """
    Get file path.
    
    Parameters
    ----------
    file_name: str
        File name.
    
    Returns
    -------
    str
        File path.
    """
    file_types = ('png', 'mp3', 'mp4')
    return [
        os.path.join(f'{file_name}.{type}') for type in file_types
    ]

def _execute_cmd(command: str) -> None:
    """
    Execute a command.
    
    Parameters
    ----------
    command: str
        Command to execute.

    Returns
    -------
    None
    """
    call(command.split(' '), stdout=DEVNULL, stderr=STDOUT)

def _pptx_pdf(pptx_file: str) -> None:
    """
    Convert pptx to pdf.
    
    Parameters
    ----------
    pptx_file: str
        Input file path to pptx file to convert.

    Returns
    -------
    None
    """
    pptx_file = os.path.abspath(pptx_file)
    powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(pptx_file)
    deck.SaveAs(os.path.abspath(f'{pptx_file[:-5]}.pdf'), 32)
    deck.Close()
    powerpoint.Quit()

def _speech_text(slide:str) -> str:
    """
    Obtain speech text from slide.

    Parameters
    ----------
    slide: str
        Slide to get speech text from.

    Returns
    -------
    str
        Speech text.
    """
    slide_text_arr = []
    presenter_notes = slide.notes_slide.notes_text_frame.text

    if len(presenter_notes) != 0 and slide.shapes.title is not None:
        title = slide.shapes.title.text
    else:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and len(shape.text) != 0:
                slide_text_arr.append(shape.text)
        title = slide_text_arr[0]

    if len(presenter_notes) == 0:
        slide_text = '. '.join(slide_text_arr)
        speech = f'{slide_text}.'
    else:
        speech = f'{title}. {presenter_notes}'
    return speech

def main():
    """Main code for script."""
    print('This script converts a powerpoint presentation (pptx) to video using text to speech with the presenter notes or slide text.')
    parser = argparse.ArgumentParser()
    required = parser.add_argument_group('required arguments')
    required.add_argument('-p', '--pptx', type=str, metavar='pptx_file', help='Input PPTX file to convert.', required = True)
    required.add_argument('-o', '--output', type=str, metavar='mp4_file', help='Output MP4 file for video.', required = True)
    args = parser.parse_args()
    pptx_video(args.pptx, args.output)
    

if __name__ == '__main__':
    main()
    #add code to resume process from temporary files
